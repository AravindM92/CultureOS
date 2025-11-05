from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Create presentation
prs = Presentation()

# Slide 1: Title
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide1.shapes.title
subtitle = slide1.placeholders[1]
title.text = "Culture OS – Thunai"
subtitle.text = "Your Team Companion\n🤖 Enhancing Team Engagement & Work-from-Office Intelligence"

# Slide 2: Problem & Solution
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide2.shapes.title
title.text = "The Challenge & Our Solution"
content = slide2.placeholders[1]
content.text = """Problem:
• Low team engagement
• Unpredictable WFO attendance
• Missed celebration moments
• Disconnected remote teams

Solution:
• AI-powered team companion
• Direct WFO intent collection
• Automated moment detection
• Seamless Teams integration"""

# Slide 3: Meet Thunai Bot
slide3 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide3.shapes.title
title.text = "Meet Thunai Bot 🤖"
content = slide3.placeholders[1]
content.text = """Personality Traits:
• 😊 Friendly - Warm, approachable conversations
• 🎉 Cheerful - Brings positive energy to teams
• 😏 Slightly Sarcastic - That fun teammate everyone enjoys

🔒 Privacy First Approach:
Every interaction requires explicit user consent. Terms & Conditions are presented before any data collection, with consent flags stored securely.

🔗 Fully Integrated with Microsoft Teams
Available in Teams playground for live demonstrations"""

# Slide 4: Architecture
slide4 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide4.shapes.title
title.text = "System Architecture"
content = slide4.placeholders[1]
content.text = """Technology Stack:
• LLM: Groq AI with Llama 3.1 8B Instant model
• Backend: FastAPI
• Database: SQLite
• Frontend: Microsoft Teams
• Bot Framework: Node.js

Flow:
Teams → Node.js Bot → Groq LLM → FastAPI → SQLite → React Dashboard"""

# Slide 5: Moments Module & Implementation
slide5 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide5.shapes.title
title.text = "Moments Module & Implementation ✨"
content = slide5.placeholders[1]
content.text = """🎯 Key Features:
• Auto-Detection: AI identifies birthdays, promotions, achievements
• Real-time Storage: Moments captured and stored instantly
• Team Celebrations: Automated notifications and greetings
• Privacy Compliant: User consent required for personal data

✅ Live Implementation:
• Bot conversations with moment detection
• Database integration with secure storage
• RESTful API with Swagger documentation
• Real-time dashboard integration"""

# Slide 6: Technical Architecture
slide6 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide6.shapes.title
title.text = "Technical Implementation 🔧"
content = slide6.placeholders[1]
content.text = """📱 Input Source:
• Microsoft Teams is the ONLY source of user interaction
• No other input channels or interfaces

🛠️ Tech Stack Details:
• Node.js Bot: Teams integration with conversation handling
• Groq LLM: Llama 3.1 8B Instant model for AI responses
• Local SQLite DB: Non-centralized, integrated database
• Python FastAPI: Thunai API for data management
• React Dashboard: Lightweight UI for analytics

⚡ Architecture Benefits:
• Lightweight and fast deployment
• Local data storage for privacy
• Scalable microservices approach"""

# Slide 7: WFO Prediction
slide7 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide7.shapes.title
title.text = "WFO Prediction Intelligence 🏢"
content = slide7.placeholders[1]
content.text = """Workflow:
1. 8 PM Check-in: Bot asks about tomorrow's office plans
2. Consent Collection: T&C acceptance for data usage
3. Intent Capture: Direct user responses stored securely
4. Smart Analytics: Real attendance patterns vs predictions

🎯 Why This Approach Works:
• Direct Intent: No guesswork - we ask users directly
• Friendly Interaction: Natural conversation, not surveys
• Real Data: Actual user intentions vs algorithmic predictions
• Privacy First: Explicit consent for every data point collected"""

# Slide 8: Additional Features
slide8 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide8.shapes.title
title.text = "Extended Capabilities 🚀"
content = slide8.placeholders[1]
content.text = """🔥 Stories [Hot Now]:
• OTT discussions • IPL polling • Chennai travel & food tips

📅 Daily Quests:
• Grab a coffee reminder • Stretch break alerts • Talk to a teammate

🎊 Seasonal Quests:
• Plan team lunches • Fun awards ceremonies • Account-level DC activities

🔮 Future Scope:
• Sentiment analysis • Emotional intelligence • Project vibe monitoring

🛡️ Admin Control & Privacy:
All features managed by account-specific Thunai coordinator with granular privacy controls and user consent management."""

# Slide 9: Why A2M TechForce
slide9 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide9.shapes.title
title.text = "Why A2M TechForce? 🏆"
content = slide9.placeholders[1]
content.text = """🎯 Use Case Coverage Leadership:
We lead in comprehensive use case coverage across the organization

🔗 Integration Ecosystem:
• Agent Architecture: Any team's application can integrate as a Thunai agent
• Unified Platform: Single bot, multiple capabilities
• Scalable Design: Easy to add new modules

🚀 Real-World Applications:
• Food Ordering: Vendor/menu lookup integration
• Carpooling: Location-based teammate matching
• Parking: Smart suggestions based on availability
• Privacy Compliant: All integrations follow T&C protocols"""

# Save presentation
prs.save('CultureOS-Thunai-Presentation.pptx')
print("PowerPoint presentation created: CultureOS-Thunai-Presentation.pptx")